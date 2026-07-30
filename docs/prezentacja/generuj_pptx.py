"""Ta sama prezentacja w formacie PowerPoint — do samodzielnego poprawiania.

Wersje HTML/PDF (generuj.py) i PPTX biorą treść z JEDNEGO zestawu stałych, więc nie
rozjadą się między sobą. Slajd ma 13,333 × 7,5 cala, czyli dokładnie 1280 × 720 px
przy 96 dpi — geometrię przenosimy więc jeden do jednego, w pikselach.

Wszystko jest natywnymi kształtami i polami tekstowymi PowerPointa (zrzuty ekranu to
jedyne obrazki), więc każdy element da się kliknąć i poprawić bez wracania do kodu.

Uruchomienie:
    python generuj_pptx.py
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from generuj import (CENNIK, CZASY, DATA, DZIALY, ETAPY_WDROZENIA, KAFELKI_LICZB,
                     KROKI_SCIEZKI, OSOBY_KONTAKT, SLAJDY, WARSTWY_SERWEROWNI, WERSJA,
                     ZRZUTY)

KATALOG = os.path.dirname(os.path.abspath(__file__))
GRANAT = RGBColor(0x1D, 0x2A, 0x4D)
TURKUS = RGBColor(0x1F, 0xC8, 0xBA)
TURKUS_CIEMNY = RGBColor(0x0F, 0x9B, 0x8E)
NIEBIESKI = RGBColor(0x25, 0x63, 0xEB)
TEKST = RGBColor(0x1E, 0x29, 0x3B)
SZARY = RGBColor(0x64, 0x74, 0x8B)
LINIA = RGBColor(0xE2, 0xE8, 0xF0)
TLO = RGBColor(0xF8, 0xFA, 0xFC)
BIALY = RGBColor(0xFF, 0xFF, 0xFF)

CZCIONKA = "Segoe UI"
MARGINES = 64          # ten sam lewy margines co w HTML


def px(v):
    """Piksele (przy 96 dpi) na jednostki PowerPointa."""
    return Emu(int(v * 914400 / 96))


def pole(slajd, x, y, w, h):
    tb = slajd.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def tekst(tf, tresc, rozmiar, kolor=TEKST, bold=False, interlinia=1.25, akapit=None):
    p = akapit if akapit is not None else (tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph())
    r = p.add_run()
    r.text = tresc
    r.font.size = Pt(rozmiar)
    r.font.color.rgb = kolor
    r.font.bold = bold
    r.font.name = CZCIONKA
    p.line_spacing = interlinia
    return p


def prostokat(slajd, x, y, w, h, wypelnienie=None, obramowanie=None, promien=True, grubosc=1):
    ksztalt = slajd.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if promien else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    if promien:
        ksztalt.adjustments[0] = 0.06
    if wypelnienie is None:
        ksztalt.fill.background()
    else:
        ksztalt.fill.solid()
        ksztalt.fill.fore_color.rgb = wypelnienie
    if obramowanie is None:
        ksztalt.line.fill.background()
    else:
        ksztalt.line.color.rgb = obramowanie
        ksztalt.line.width = Pt(grubosc)
    ksztalt.shadow.inherit = False
    ksztalt.text_frame.word_wrap = True
    return ksztalt


def naglowek(slajd, tytul, numer):
    tf = pole(slajd, MARGINES, 48, 1100, 60)
    tekst(tf, tytul, 30, GRANAT, bold=True, interlinia=1.1)
    kreska = prostokat(slajd, MARGINES, 118, 64, 5, TURKUS, promien=False)
    kreska.line.fill.background()
    stopka = pole(slajd, MARGINES, 686, 500, 20)
    tekst(stopka, "ZCO Document Management", 9.5, RGBColor(0x94, 0xA3, 0xB8))
    nr = pole(slajd, 1180, 686, 60, 20)
    p = tekst(nr, str(numer), 9.5, RGBColor(0x94, 0xA3, 0xB8))
    p.alignment = PP_ALIGN.RIGHT


def puenta(slajd, tresc):
    pas = prostokat(slajd, MARGINES, 596, 1152, 62, GRANAT)
    akcent = prostokat(slajd, MARGINES, 596, 6, 62, TURKUS, promien=False)
    akcent.line.fill.background()
    tf = pas.text_frame
    tf.margin_left, tf.margin_top = px(22), px(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    akapit = tekst(tf, tresc, 16.5, BIALY, bold=True)
    akapit.alignment = PP_ALIGN.LEFT


def punktory(slajd, x, y, w, pozycje, rozmiar=17):
    """Lista punktowana: kropka jako osobny kształt, żeby zachować kolor marki."""
    gora = y
    for tresc in pozycje:
        kropka = slajd.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(gora + 8), px(9), px(9))
        kropka.fill.solid()
        kropka.fill.fore_color.rgb = TURKUS_CIEMNY
        kropka.line.fill.background()
        kropka.shadow.inherit = False
        tf = pole(slajd, x + 24, gora, w - 24, 30)
        tekst(tf, tresc, rozmiar, TEKST, interlinia=1.3)
        znakow_w_wierszu = max(20, int((w - 24) / (rozmiar * 0.62)))
        wierszy = max(1, -(-len(tresc) // znakow_w_wierszu))
        gora += 30 * wierszy + 16
    return gora


def bez_html(t):
    import re
    return re.sub(r"<[^>]+>", "", t).replace("&nbsp;", " ")


# ---------------------------------------------------------------- slajdy

def slajd_okladka(prs, dane):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tlo = prostokat(s, 0, 0, 1280, 720, GRANAT, promien=False)
    tlo.line.fill.background()
    kreska = prostokat(s, MARGINES, 246, 88, 6, TURKUS, promien=False)
    kreska.line.fill.background()
    tf = pole(s, MARGINES, 282, 1000, 90)
    tekst(tf, dane["tytul"], 48, BIALY, bold=True, interlinia=1.05)
    tf2 = pole(s, MARGINES, 392, 860, 80)
    for linia in bez_html(dane["podtytul"]).split("  ") if False else bez_html(
            dane["podtytul"].replace("<br>", "\n")).split("\n"):
        tekst(tf2, linia, 19, TURKUS, interlinia=1.45)
    tf3 = pole(s, MARGINES, 622, 700, 50)
    tekst(tf3, "Polmedi Group sp. z o.o. · Poznań", 12, RGBColor(0xE2, 0xE8, 0xF0), bold=True)
    tekst(tf3, f"Wersja aplikacji {WERSJA} · {DATA}", 12, RGBColor(0x94, 0xA3, 0xB8))


def graf_serwerownia(s, y):
    ramka = prostokat(s, MARGINES, y, 800, 190, GRANAT)
    tf = pole(s, MARGINES + 20, y + 18, 500, 20)
    tekst(tf, "SZPITAL · WASZA SEROWNIA".replace("SEROWNIA", "SERWEROWNIA"), 11, TURKUS, bold=True)
    szer = (800 - 40 - 2 * 12) / 3
    for i, (_, nazwa, opis) in enumerate(WARSTWY_SERWEROWNI):
        x = MARGINES + 20 + i * (szer + 12)
        k = prostokat(s, x, y + 52, szer, 118, RGBColor(0x2A, 0x3A, 0x63))
        k.line.fill.background()
        tfk = pole(s, x + 16, y + 70, szer - 32, 90)
        tekst(tfk, nazwa, 13.5, BIALY, bold=True)
        tekst(tfk, opis, 10.5, RGBColor(0xCB, 0xD5, 0xE1), interlinia=1.3)
    # przerywana kreska w kolorze marki
    linia = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(MARGINES + 852), px(y + 24), px(0), px(96))
    linia.line.color.rgb = TURKUS_CIEMNY
    linia.line.width = Pt(3)
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    linia.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    linia.shadow.inherit = False
    tfm = pole(s, MARGINES + 800, y + 132, 104, 50)
    pm = tekst(tfm, "brak połączenia z usługami zewnętrznymi", 9.5, TURKUS_CIEMNY, bold=True, interlinia=1.25)
    pm.alignment = PP_ALIGN.CENTER
    # kafelek Internet — pełne kolory, czytelny
    net = prostokat(s, MARGINES + 906, y, 246, 190, TLO, RGBColor(0xCB, 0xD5, 0xE1), grubosc=1.5)
    tfn = pole(s, MARGINES + 922, y + 46, 214, 120)
    pn = tekst(tfn, "Internet", 14, GRANAT, bold=True)
    pn.alignment = PP_ALIGN.CENTER
    pn2 = tekst(tfn, "nie wychodzi tu żaden dokument, fragment ani pytanie", 10.5, TEKST, interlinia=1.35)
    pn2.alignment = PP_ALIGN.CENTER
    return y + 210


def graf_kroki(s, y):
    szer = (1152 - 3 * 26) / 4
    for i, (nr, nazwa, opis) in enumerate(KROKI_SCIEZKI):
        x = MARGINES + i * (szer + 26)
        prostokat(s, x, y, szer, 130, TLO, LINIA)
        kolo = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x + 16), px(y + 14), px(26), px(26))
        kolo.fill.solid()
        kolo.fill.fore_color.rgb = TURKUS_CIEMNY
        kolo.line.fill.background()
        kolo.shadow.inherit = False
        tfk = kolo.text_frame
        tfk.word_wrap = False
        tfk.margin_left = tfk.margin_right = tfk.margin_top = tfk.margin_bottom = 0
        pk = tekst(tfk, nr, 11, BIALY, bold=True, interlinia=1.0)
        pk.alignment = PP_ALIGN.CENTER
        tf = pole(s, x + 16, y + 50, szer - 32, 70)
        tekst(tf, nazwa, 14, GRANAT, bold=True)
        tekst(tf, opis, 10.5, SZARY, interlinia=1.3)
        if i < 3:
            tfs = pole(s, x + szer + 4, y + 52, 20, 24)
            ps = tekst(tfs, "→", 15, TURKUS_CIEMNY)
            ps.alignment = PP_ALIGN.CENTER
    return y + 150


def graf_liczby_pptx(s, y):
    szer = (1152 - 2 * 18) / 3
    for i, (duza, opis, nota) in enumerate(KAFELKI_LICZB):
        x = MARGINES + i * (szer + 18)
        prostokat(s, x, y, szer, 130, TLO, LINIA)
        tf = pole(s, x + 24, y + 22, szer - 48, 100)
        tekst(tf, duza, 34, GRANAT, bold=True, interlinia=1.0)
        tekst(tf, opis, 13.5, TEKST, interlinia=1.3)
        tekst(tf, nota, 10, SZARY, interlinia=1.3)
    y += 152
    tf = pole(s, MARGINES, y, 400, 22)
    tekst(tf, "Ile to trwa w praktyce", 11.5, GRANAT, bold=True)
    y += 30
    TOR = 620
    for nazwa, sekundy, etykieta in CZASY:
        tfn = pole(s, MARGINES, y - 2, 300, 24)
        tekst(tfn, nazwa, 12.5, TEKST)
        szerokosc = max(6, TOR * sekundy / 60)      # długość wprost proporcjonalna do czasu
        pasek = prostokat(s, MARGINES + 310, y, szerokosc, 15,
                          TURKUS_CIEMNY if sekundy < 60 else NIEBIESKI, promien=False)
        pasek.line.fill.background()
        tfw = pole(s, MARGINES + 310 + TOR + 14, y - 2, 200, 24)
        tekst(tfw, etykieta, 12.5, TEKST, bold=True)
        y += 30
    tf = pole(s, MARGINES, y + 4, 1100, 30)
    tekst(tf, "Długość słupka odpowiada czasowi. Pomiary z działającego wdrożenia demonstracyjnego; "
              "przygotowanie dokumentu odbywa się raz, w tle, przy wgraniu pliku.", 9.5, SZARY, interlinia=1.35)
    return y + 40


def graf_etapy(s, y):
    szer = (1152 - 2 * 14) / 3
    for i, (tydzien, nazwa, opis) in enumerate(ETAPY_WDROZENIA):
        x = MARGINES + i * (szer + 14)
        pion = prostokat(s, x, y, 3, 104, TURKUS, promien=False)
        pion.line.fill.background()
        kolo = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x - 6), px(y + 4), px(15), px(15))
        kolo.fill.solid()
        kolo.fill.fore_color.rgb = TURKUS_CIEMNY
        kolo.line.color.rgb = BIALY
        kolo.line.width = Pt(2)
        kolo.shadow.inherit = False
        tf = pole(s, x + 16, y + 2, szer - 24, 100)
        tekst(tf, tydzien.upper(), 10, TURKUS_CIEMNY, bold=True)
        tekst(tf, nazwa, 15, GRANAT, bold=True)
        tekst(tf, opis, 11.5, SZARY, interlinia=1.35)
    return y + 124


def graf_dzialy(s, y):
    szer = (1152 - 2 * 14) / 3
    for i, (nazwa, opis) in enumerate(DZIALY):
        x = MARGINES + (i % 3) * (szer + 14)
        gora = y + (i // 3) * 96
        prostokat(s, x, gora, szer, 82, TLO, LINIA)
        tf = pole(s, x + 20, gora + 16, szer - 40, 60)
        tekst(tf, nazwa, 14.5, GRANAT, bold=True)
        tekst(tf, opis, 11, SZARY, interlinia=1.3)
    return y + 192


def graf_cennik_pptx(s, y):
    tabela = s.shapes.add_table(5, 2, px(MARGINES), px(y), px(1152), px(230)).table
    tabela.columns[0].width, tabela.columns[1].width = px(880), px(272)
    wiersze = [("POZYCJA", "NETTO")] + CENNIK + [("Rok pierwszy łącznie", "57 000 zł")]
    for i, (lewa, prawa) in enumerate(wiersze):
        for j, wartosc in enumerate((lewa, prawa)):
            kom = tabela.cell(i, j)
            kom.text = wartosc
            kom.fill.solid()
            kom.fill.fore_color.rgb = BIALY
            kom.margin_left = kom.margin_right = px(10)
            p = kom.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j == 1 else PP_ALIGN.LEFT
            r = p.runs[0]
            r.font.name = CZCIONKA
            if i == 0:
                r.font.size, r.font.bold, r.font.color.rgb = Pt(10), True, SZARY
            elif i == len(wiersze) - 1:
                r.font.size, r.font.bold, r.font.color.rgb = Pt(16), True, GRANAT
            else:
                r.font.size, r.font.color.rgb = Pt(14), TEKST
    y += 244
    noty = ["Bez opłat za użytkownika i bez opłat za liczbę dokumentów.",
            "Sprzęt zostaje Wasz — nie płacicie abonamentu za dostęp do własnych danych.",
            "Trzy lata użytkowania: 81 000 zł netto."]
    szer = (1152 - 2 * 28) / 3
    for i, nota in enumerate(noty):
        x = MARGINES + i * (szer + 28)
        pion = prostokat(s, x, y, 3, 52, TURKUS, promien=False)
        pion.line.fill.background()
        tf = pole(s, x + 12, y, szer - 12, 52)
        tekst(tf, nota, 11.5, TEKST, interlinia=1.35)
    return y + 60


def graf_kontakt_pptx(s, y):
    tf = pole(s, MARGINES, y, 600, 26)
    tekst(tf, "Czekają na Was:", 17, TEKST)
    y += 38
    szer = (1152 - 28) / 2
    for i, (imie, rola, telefon) in enumerate(OSOBY_KONTAKT):
        x = MARGINES + i * (szer + 28)
        prostokat(s, x, y, szer, 132, TLO, LINIA)
        pion = prostokat(s, x, y, 6, 132, TURKUS_CIEMNY, promien=False)
        pion.line.fill.background()
        kolo = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x + 34), px(y + 30), px(72), px(72))
        kolo.fill.solid()
        kolo.fill.fore_color.rgb = GRANAT
        kolo.line.fill.background()
        kolo.shadow.inherit = False
        czlony = imie.split()
        tfi = kolo.text_frame
        tfi.word_wrap = False
        tfi.margin_left = tfi.margin_right = tfi.margin_top = tfi.margin_bottom = 0
        pk = tekst(tfi, (czlony[0][0] + czlony[-1][0]).upper(), 20, BIALY, bold=True, interlinia=1.0)
        pk.alignment = PP_ALIGN.CENTER
        tf = pole(s, x + 126, y + 28, szer - 150, 90)
        tekst(tf, imie, 19, GRANAT, bold=True, interlinia=1.15)
        tekst(tf, rola, 12.5, SZARY)
        tekst(tf, f"tel. {telefon}", 18, TURKUS_CIEMNY, bold=True, interlinia=1.3)
    y += 152
    kreska = prostokat(s, MARGINES, y, 1152, 1, LINIA, promien=False)
    kreska.line.fill.background()
    tf = pole(s, MARGINES, y + 16, 800, 26)
    tekst(tf, "Polmedi Group sp. z o.o. · Poznań · polmedi.com", 13.5, GRANAT, bold=True)
    return y + 50


GRAFIKI = {
    "Bezpieczeństwo i prywatność": graf_serwerownia,
    "Skąd pewność, że odpowiedź jest prawdziwa": graf_kroki,
    "Pojemność i szybkość": graf_liczby_pptx,
    "Nasza rola: wdrożenie i wsparcie": graf_etapy,
    "Gdzie to pracuje": graf_dzialy,
    "Warunki i następny krok": graf_cennik_pptx,
    "Zapraszamy do kontaktu": graf_kontakt_pptx,
}


def punkty_ze_slajdu(dane):
    """Wyciągnij same punkty listy z HTML-owej treści slajdu."""
    import re
    return [bez_html(x) for x in re.findall(r"<li>(.*?)</li>", dane.get("tresc", ""), re.S)]


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = px(1280), px(720)

    for numer, dane in enumerate(SLAJDY, start=1):
        if dane.get("typ") == "okladka":
            slajd_okladka(prs, dane)
            continue

        s = prs.slides.add_slide(prs.slide_layouts[6])
        naglowek(s, dane["tytul"], numer)
        y = 170

        rysuj = GRAFIKI.get(dane["tytul"])
        if rysuj:
            y = rysuj(s, y) + 12

        lista = punkty_ze_slajdu(dane)
        if lista:
            szerokosc = 620 if dane.get("obraz") else 1152
            punktory(s, MARGINES, y, szerokosc, lista)

        if dane.get("obraz"):
            sciezka = os.path.join(ZRZUTY, dane["obraz"])
            if os.path.exists(sciezka):
                s.shapes.add_picture(sciezka, px(700), px(176), width=px(516))

        if dane.get("puenta"):
            puenta(s, dane["puenta"])

        if dane.get("notatka"):
            s.notes_slide.notes_text_frame.text = dane["notatka"]

    cel = os.path.join(KATALOG, "ZCO-DM-prezentacja.pptx")
    prs.save(cel)
    print(f"ZCO-DM-prezentacja.pptx: {len(prs.slides.__iter__.__self__._sldIdLst)} slajdów, "
          f"{os.path.getsize(cel)//1024} KB")


if __name__ == "__main__":
    main()
