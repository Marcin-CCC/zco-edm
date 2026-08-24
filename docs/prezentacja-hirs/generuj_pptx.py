"""HiRS — ta sama prezentacja w formacie PowerPoint, do samodzielnego poprawiania.

Treść bierze z tych samych stałych co wersja HTML/PDF (`generuj.py`), więc obie nie mogą
się rozjechać. Wszystko jest natywnymi kształtami i polami tekstowymi; jedyne obrazy to
logo i dwie makiety ekranów (renderowane przez `makiety_png.py` z tego samego źródła co
w HTML).

Slajd 13,333 × 7,5 cala = 1280 × 720 px przy 96 dpi, więc geometrię przenosimy jeden do
jednego, w pikselach.

    python makiety_png.py && python generuj_pptx.py
"""
import os
import re

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt

from generuj import (CENNIK, CZASY, DATA, DZIALY, ETAPY_WDROZENIA, KAFELKI_LICZB, KATALOG,
                     KROKI_SCIEZKI, LOGO_KONTRA_PNG, LOGO_PNG, NAZWA, OSOBY_KONTAKT,
                     PODTYTUL_NAZWY, SLAJDY, WARSTWY_SERWEROWNI, WYKONAWCA)
from generuj import LINIA as _LINIA, NIEBIESKI as _NIEBIESKI, SZARY as _SZARY
from generuj import TEKST as _TEKST, TLO as _TLO, TURKUS as _TURKUS
from generuj import TURKUS_CIEMNY as _TURKUS_CIEMNY


def _rgb(hex_kolor: str) -> RGBColor:
    """„#1d2a4d" → RGBColor. Kolory bierzemy z `generuj`, a nie przepisujemy tutaj.

    Wcześniej obie wersje trzymały własną kopię palety i przy zmianie marki
    rozjechały się bez ostrzeżenia: HTML był granatowy, a PowerPoint nadal
    niebieski — wyglądało to na dwie różne prezentacje.
    """
    s = hex_kolor.lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _rgb_jasniej(hex_kolor: str, krycie_bieli: float) -> RGBColor:
    """Kolor zmieszany z bielą — zamiennik przezroczystości, której PowerPoint
    nie ma dla wypełnień kształtów."""
    s = hex_kolor.lstrip("#")
    skladowe = (int(s[i:i + 2], 16) for i in (0, 2, 4))
    return RGBColor(*(int(round(k + (255 - k) * krycie_bieli)) for k in skladowe))


NIEBIESKI = _rgb(_NIEBIESKI)
TURKUS = _rgb(_TURKUS)
TURKUS_CIEMNY = _rgb(_TURKUS_CIEMNY)
TEKST = _rgb(_TEKST)
SZARY = _rgb(_SZARY)
LINIA = _rgb(_LINIA)
TLO = _rgb(_TLO)
BIALY = RGBColor(0xFF, 0xFF, 0xFF)

CZCIONKA = "Segoe UI"
MARGINES = 64


def px(v):
    return Emu(int(v * 914400 / 96))


def pole(slajd, x, y, w, h):
    tb = slajd.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def tekst(tf, tresc, rozmiar, kolor=TEKST, bold=False, interlinia=1.25):
    p = tf.paragraphs[0] if not tf.paragraphs[0].runs else tf.add_paragraph()
    r = p.add_run()
    r.text = tresc
    r.font.size = Pt(rozmiar)
    r.font.color.rgb = kolor
    r.font.bold = bold
    r.font.name = CZCIONKA
    p.line_spacing = interlinia
    return p


def prostokat(slajd, x, y, w, h, wypelnienie=None, obramowanie=None, promien=True, grubosc=1):
    k = slajd.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if promien else MSO_SHAPE.RECTANGLE,
        px(x), px(y), px(w), px(h))
    if promien:
        k.adjustments[0] = 0.06
    if wypelnienie is None:
        k.fill.background()
    else:
        k.fill.solid()
        k.fill.fore_color.rgb = wypelnienie
    if obramowanie is None:
        k.line.fill.background()
    else:
        k.line.color.rgb = obramowanie
        k.line.width = Pt(grubosc)
    k.shadow.inherit = False
    k.text_frame.word_wrap = True
    return k


def pasek_gradientu(slajd, x, y, w, h):
    """Kreska w przejściu niebiesko-turkusowym — znak rozpoznawczy marki."""
    k = prostokat(slajd, x, y, w, h, NIEBIESKI, promien=False)
    k.fill.gradient()
    k.fill.gradient_angle = 0
    k.fill.gradient_stops[0].color.rgb = NIEBIESKI
    k.fill.gradient_stops[0].position = 0.0
    k.fill.gradient_stops[1].color.rgb = TURKUS
    k.fill.gradient_stops[1].position = 1.0
    k.line.fill.background()
    return k


def bez_html(t):
    return re.sub(r"<[^>]+>", "", t).replace("&nbsp;", " ")


def naglowek(slajd, tytul, numer, logo):
    tf = pole(slajd, MARGINES, 48, 1000, 60)
    tekst(tf, tytul, 30, NIEBIESKI, bold=True, interlinia=1.1)
    pasek_gradientu(slajd, MARGINES, 118, 96, 5)
    if logo:
        slajd.shapes.add_picture(logo, px(1108), px(44), width=px(108))
    stopka = pole(slajd, MARGINES, 686, 500, 20)
    tekst(stopka, NAZWA, 9.5, RGBColor(0x9A, 0xA5, 0xA5))
    nr = pole(slajd, 1180, 686, 60, 20)
    tekst(nr, str(numer), 9.5, RGBColor(0x9A, 0xA5, 0xA5)).alignment = PP_ALIGN.RIGHT


def puenta(slajd, tresc):
    pas = prostokat(slajd, MARGINES, 596, 1152, 62, NIEBIESKI)
    pasek_gradientu(slajd, MARGINES, 596, 6, 62)
    tf = pas.text_frame
    tf.margin_left, tf.margin_top = px(22), px(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tekst(tf, tresc, 16.5, BIALY, bold=True).alignment = PP_ALIGN.LEFT


def punktory(slajd, x, y, w, pozycje, rozmiar=17):
    gora = y
    for tresc in pozycje:
        kropka = slajd.shapes.add_shape(MSO_SHAPE.OVAL, px(x), px(gora + 8), px(9), px(9))
        kropka.fill.solid()
        kropka.fill.fore_color.rgb = TURKUS_CIEMNY
        kropka.line.fill.background()
        kropka.shadow.inherit = False
        tf = pole(slajd, x + 24, gora, w - 24, 30)
        tekst(tf, tresc, rozmiar, TEKST, interlinia=1.3)
        znakow = max(20, int((w - 24) / (rozmiar * 0.62)))
        gora += 30 * max(1, -(-len(tresc) // znakow)) + 16
    return gora


# ---------------------------------------------------------------- slajdy

def slajd_okladka(prs, dane, logo, logo_kontra):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tlo = prostokat(s, 0, 0, 1280, 720, NIEBIESKI, promien=False)
    tlo.line.fill.background()
    pasek_gradientu(s, MARGINES, 246, 120, 6)
    tf = pole(s, MARGINES, 286, 1000, 100)
    tekst(tf, dane["tytul"], 58, BIALY, bold=True, interlinia=1.02)
    tf2 = pole(s, MARGINES, 372, 900, 30)
    tekst(tf2, PODTYTUL_NAZWY, 15, RGBColor(0xC7, 0xF3, 0xF3))
    tf3 = pole(s, MARGINES, 414, 880, 80)
    for linia in bez_html(dane["podtytul"].replace("<br>", "\n")).split("\n"):
        tekst(tf3, linia, 19, BIALY, interlinia=1.45)
    tf4 = pole(s, MARGINES, 626, 700, 50)
    tekst(tf4, WYKONAWCA, 12, RGBColor(0xEA, 0xF0, 0xFF), bold=True)
    tekst(tf4, DATA, 12, RGBColor(0xB9, 0xC6, 0xEA))
    # Znak w kontrze kładziemy wprost na granacie. Wcześniej stał tu biały prostokąt
    # z logiem pełnokolorowym — na ciemnym tle wyglądał jak naklejka.
    if logo_kontra:
        szer, wys = 207, 46         # proporcja znaku 1999 × 444
        s.shapes.add_picture(logo_kontra, px(1216 - szer), px(676 - wys), width=px(szer))
    elif logo:
        biale = prostokat(s, 1064, 636, 152, 48, BIALY)
        biale.line.fill.background()
        s.shapes.add_picture(logo, px(1080), px(646), width=px(120))


def graf_serwerownia(s, y):
    prostokat(s, MARGINES, y, 800, 190, NIEBIESKI)
    tekst(pole(s, MARGINES + 20, y + 18, 500, 20), "WASZA SERWEROWNIA", 11, TURKUS, bold=True)
    szer = (800 - 40 - 2 * 12) / 3
    for i, (nazwa, opis) in enumerate(WARSTWY_SERWEROWNI):
        x = MARGINES + 20 + i * (szer + 12)
        # HTML kładzie tu białą warstwę o krycu 10% na kolorze wiodącym; PowerPoint
        # nie zna przezroczystości wypełnienia, więc podajemy wynik zmieszania wprost.
        k = prostokat(s, x, y + 52, szer, 118, _rgb_jasniej(_NIEBIESKI, 0.10))
        k.line.fill.background()
        tf = pole(s, x + 16, y + 72, szer - 32, 90)
        tekst(tf, nazwa, 13.5, BIALY, bold=True)
        tekst(tf, opis, 10.5, RGBColor(0xDB, 0xE4, 0xF7), interlinia=1.3)
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    linia = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, px(MARGINES + 852), px(y + 24), px(0), px(96))
    linia.line.color.rgb = TURKUS_CIEMNY
    linia.line.width = Pt(3)
    linia.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    linia.shadow.inherit = False
    tf = pole(s, MARGINES + 800, y + 132, 104, 50)
    tekst(tf, "brak połączenia z usługami zewnętrznymi", 9.5, TURKUS_CIEMNY, bold=True,
          interlinia=1.25).alignment = PP_ALIGN.CENTER
    prostokat(s, MARGINES + 906, y, 246, 190, TLO, RGBColor(0xCB, 0xD5, 0xE1), grubosc=1.5)
    tf = pole(s, MARGINES + 922, y + 56, 214, 110)
    tekst(tf, "Internet", 14, NIEBIESKI, bold=True).alignment = PP_ALIGN.CENTER
    tekst(tf, "nie wychodzi tu żaden dokument, fragment ani pytanie", 10.5, TEKST,
          interlinia=1.35).alignment = PP_ALIGN.CENTER
    return y + 210


def graf_kroki(s, y):
    szer = (1152 - 3 * 26) / 4
    for i, (nr, nazwa, opis) in enumerate(KROKI_SCIEZKI):
        x = MARGINES + i * (szer + 26)
        prostokat(s, x, y, szer, 130, TLO, LINIA)
        kolo = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x + 16), px(y + 14), px(26), px(26))
        kolo.fill.solid()
        kolo.fill.fore_color.rgb = TURKUS_CIEMNY
        kolo.line.fill.background()
        kolo.shadow.inherit = False
        tfk = kolo.text_frame
        tfk.word_wrap = False
        tfk.margin_left = tfk.margin_right = tfk.margin_top = tfk.margin_bottom = 0
        tekst(tfk, nr, 11, BIALY, bold=True, interlinia=1.0).alignment = PP_ALIGN.CENTER
        tf = pole(s, x + 16, y + 50, szer - 32, 70)
        tekst(tf, nazwa, 14, NIEBIESKI, bold=True)
        tekst(tf, opis, 10.5, SZARY, interlinia=1.3)
        if i < 3:
            tf = pole(s, x + szer + 4, y + 52, 20, 24)
            tekst(tf, "→", 15, TURKUS_CIEMNY).alignment = PP_ALIGN.CENTER
    return y + 150


def graf_liczby(s, y):
    szer = (1152 - 2 * 18) / 3
    for i, (duza, opis, nota) in enumerate(KAFELKI_LICZB):
        x = MARGINES + i * (szer + 18)
        prostokat(s, x, y, szer, 130, TLO, LINIA)
        tf = pole(s, x + 24, y + 22, szer - 48, 100)
        tekst(tf, duza, 34, NIEBIESKI, bold=True, interlinia=1.0)
        tekst(tf, opis, 13.5, TEKST, interlinia=1.3)
        tekst(tf, nota, 10, SZARY, interlinia=1.3)
    y += 152
    tekst(pole(s, MARGINES, y, 400, 22), "Ile to trwa w praktyce", 11.5, NIEBIESKI, bold=True)
    y += 30
    TOR = 620
    for nazwa, sekundy, etykieta in CZASY:
        tekst(pole(s, MARGINES, y - 2, 300, 24), nazwa, 12.5, TEKST)
        pasek = prostokat(s, MARGINES + 310, y, max(6, TOR * sekundy / 60), 15,
                          TURKUS_CIEMNY if sekundy < 60 else NIEBIESKI, promien=False)
        pasek.line.fill.background()
        tekst(pole(s, MARGINES + 310 + TOR + 14, y - 2, 200, 24), etykieta, 12.5, TEKST, bold=True)
        y += 30
    tf = pole(s, MARGINES, y + 4, 1100, 30)
    tekst(tf, "Długość słupka odpowiada czasowi. Pomiary z działającego wdrożenia; przygotowanie "
              "dokumentu odbywa się raz, w tle, przy wgraniu pliku.", 9.5, SZARY, interlinia=1.35)
    return y + 40


def graf_etapy(s, y):
    szer = (1152 - 2 * 14) / 3
    for i, (tydzien, nazwa, opis) in enumerate(ETAPY_WDROZENIA):
        x = MARGINES + i * (szer + 14)
        pion = prostokat(s, x, y, 3, 120, TURKUS, promien=False)
        pion.line.fill.background()
        kolo = s.shapes.add_shape(MSO_SHAPE.OVAL, px(x - 6), px(y + 4), px(15), px(15))
        kolo.fill.solid()
        kolo.fill.fore_color.rgb = TURKUS_CIEMNY
        kolo.line.color.rgb = BIALY
        kolo.line.width = Pt(2)
        kolo.shadow.inherit = False
        tf = pole(s, x + 16, y + 2, szer - 24, 100)
        tekst(tf, tydzien.upper(), 10, TURKUS_CIEMNY, bold=True)
        tekst(tf, nazwa, 15, NIEBIESKI, bold=True)
        tekst(tf, opis, 11.5, SZARY, interlinia=1.35)
    return y + 124


def graf_dzialy(s, y):
    szer = (1152 - 2 * 14) / 3
    for i, (nazwa, opis) in enumerate(DZIALY):
        x = MARGINES + (i % 3) * (szer + 14)
        gora = y + (i // 3) * 106
        prostokat(s, x, gora, szer, 94, TLO, LINIA)
        tf = pole(s, x + 20, gora + 16, szer - 40, 72)
        tekst(tf, nazwa, 14.5, NIEBIESKI, bold=True)
        tekst(tf, opis, 11, SZARY, interlinia=1.3)
    return y + 212


def graf_cennik(s, y):
    tabela = s.shapes.add_table(5, 2, px(MARGINES), px(y), px(1152), px(230)).table
    tabela.columns[0].width, tabela.columns[1].width = px(880), px(272)
    wiersze = [("POZYCJA", "NETTO")] + CENNIK + [("Rok pierwszy łącznie", "57 000 zł")]
    for i, (lewa, prawa) in enumerate(wiersze):
        for j, wartosc in enumerate((lewa, prawa)):
            kom = tabela.cell(i, j)
            kom.text = wartosc
            kom.fill.solid()
            kom.fill.fore_color.rgb = BIALY
            kom.margin_left = kom.margin_right = px(10)
            p = kom.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j == 1 else PP_ALIGN.LEFT
            r = p.runs[0]
            r.font.name = CZCIONKA
            if i == 0:
                r.font.size, r.font.bold, r.font.color.rgb = Pt(10), True, SZARY
            elif i == len(wiersze) - 1:
                r.font.size, r.font.bold, r.font.color.rgb = Pt(16), True, NIEBIESKI
            else:
                r.font.size, r.font.color.rgb = Pt(14), TEKST
    y += 244
    noty = ["Bez opłat za użytkownika i bez opłat za liczbę dokumentów.",
            "Sprzęt zostaje Wasz — nie płacicie abonamentu za dostęp do własnych danych.",
            "Trzy lata użytkowania: 81 000 zł netto."]
    szer = (1152 - 2 * 28) / 3
    for i, nota in enumerate(noty):
        x = MARGINES + i * (szer + 28)
        pion = prostokat(s, x, y, 3, 52, TURKUS, promien=False)
        pion.line.fill.background()
        tekst(pole(s, x + 12, y, szer - 12, 52), nota, 11.5, TEKST, interlinia=1.35)
    return y + 60


def graf_cennik_kontakt(s, y):
    """Cennik i kontakt na jednym slajdzie — odpowiednik `graf_kontakt_pasek` z HTML.

    Osoba mieści się w jednym wierszu, bo pod cennikiem zostaje tylko tyle miejsca.
    Duże wizytówki (obwódka, inicjały w kole) nachodziłyby tu na pasek puenty — sprawdzone przy pierwszym złożeniu.
    """
    y = graf_cennik(s, y) + 16      # notki cennika kończą się tuż nad paskiem
    szer = (1152 - 24) / 2
    for i, (imie, rola, telefon) in enumerate(OSOBY_KONTAKT):
        x = MARGINES + i * (szer + 24)
        prostokat(s, x, y, szer, 46, TLO, LINIA)
        tekst(pole(s, x + 16, y + 11, 190, 24), imie, 14, NIEBIESKI, bold=True)
        tekst(pole(s, x + 200, y + 14, 190, 22), rola, 10.5, SZARY)
        t = pole(s, x + szer - 200, y + 11, 184, 24)
        tekst(t, f"tel. {telefon}", 13.5, TURKUS_CIEMNY, bold=True).alignment = PP_ALIGN.RIGHT
    return y + 58


GRAFIKI = {
    "Bezpieczeństwo i prywatność": graf_serwerownia,
    "Skąd pewność, że odpowiedź jest prawdziwa": graf_kroki,
    "Pojemność i szybkość": graf_liczby,
    "Nasza rola: wdrożenie i wsparcie": graf_etapy,
    "Gdzie to pracuje": graf_dzialy,
    # Wydanie HiRS ma dziesięć slajdów: warunki i kontakt stoją razem na ostatnim.
    "Warunki, następny krok i kontakt": graf_cennik_kontakt,
}
MAKIETY = {"Co dostajecie": "makieta-pliki.png", "Łatwość obsługi": "makieta-czat.png"}


def main():
    logo = os.path.join(KATALOG, LOGO_PNG)
    logo = logo if os.path.exists(logo) else None
    # Wersję w kontrze rasteryzuje `makiety_png.py`; bez niej okładka
    # wraca do loga pełnokolorowego na białym podkładzie.
    logo_kontra = os.path.join(KATALOG, LOGO_KONTRA_PNG)
    logo_kontra = logo_kontra if os.path.exists(logo_kontra) else None

    prs = Presentation()
    prs.slide_width, prs.slide_height = px(1280), px(720)

    for numer, dane in enumerate(SLAJDY, start=1):
        if dane.get("typ") == "okladka":
            slajd_okladka(prs, dane, logo, logo_kontra)
            continue

        s = prs.slides.add_slide(prs.slide_layouts[6])
        naglowek(s, dane["tytul"], numer, logo)
        y = 170

        rysuj = GRAFIKI.get(dane["tytul"])
        if rysuj:
            y = rysuj(s, y) + 12

        makieta = MAKIETY.get(dane["tytul"])
        if makieta:
            sciezka = os.path.join(KATALOG, makieta)
            if os.path.exists(sciezka):
                s.shapes.add_picture(sciezka, px(700), px(180), width=px(516))

        lista = [bez_html(x) for x in re.findall(r"<li>(.*?)</li>", dane.get("tresc", ""), re.S)]
        if lista:
            punktory(s, MARGINES, y, 620 if makieta else 1152, lista)

        if dane.get("puenta"):
            puenta(s, dane["puenta"])
        if dane.get("notatka"):
            s.notes_slide.notes_text_frame.text = dane["notatka"]

    cel = os.path.join(KATALOG, "HiRS-prezentacja.pptx")
    prs.save(cel)
    print(f"HiRS-prezentacja.pptx: {len(SLAJDY)} slajdów, {os.path.getsize(cel)//1024} KB")


if __name__ == "__main__":
    main()
